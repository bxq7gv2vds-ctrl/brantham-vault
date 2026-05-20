"""Generate Brantham Teaser — exact match with note-strategique style.
Pure B&W, CMU Serif only, booktabs, no color, no navy, no AI pizzazz.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

FONT  = "CMU Serif"
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x55, 0x55, 0x55)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUTPUT  = os.path.expanduser("~/vault/08_Exports/Brantham-Teaser-Template.pptx")


def bg_white(s): 
    f = s.background.fill; f.solid(); f.fore_color.rgb = WHITE

def tb(s, l, t, w, h, text="", size=14, bold=False, italic=False, color=BLACK,
       align=PP_ALIGN.LEFT, spacing=1.35):
    tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.name = FONT; p.font.size = Pt(size)
    p.font.bold = bold; p.font.italic = italic; p.font.color.rgb = color
    p.alignment = align; p.space_after = Pt(0); p.space_before = Pt(0)
    p.line_spacing = Pt(size * spacing)
    return tf

def add_p(tf, text, size=14, bold=False, italic=False, color=BLACK,
          align=PP_ALIGN.LEFT, before=0, after=0):
    p = tf.add_paragraph(); p.text = text
    p.font.name = FONT; p.font.size = Pt(size)
    p.font.bold = bold; p.font.italic = italic; p.font.color.rgb = color
    p.alignment = align; p.space_before = Pt(before); p.space_after = Pt(after)
    return p

def rule(s, y, l=1.0, w=11.3, height=0.5):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(y), Inches(w), Pt(height))
    r.fill.solid(); r.fill.fore_color.rgb = BLACK; r.line.fill.background()

def page_num(s):
    """Reserve space for number — user adds manually for now"""
    pass


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER (exact match note stratégique)
# ═══════════════════════════════════════════════════════════════
def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    tb(s, 2.0, 1.6, 9.3, 1.0, "Note d'analyse stratégique", size=28, bold=True,
       color=BLACK, align=PP_ALIGN.CENTER)

    tf = tb(s, 1.5, 2.8, 10.3, 2.5, "",
            size=16, color=BLACK, align=PP_ALIGN.CENTER)
    # subtitle (edit these per deal)
    add_p(tf, "Thèse de reprise et scénarios de création de valeur", 16, align=PP_ALIGN.CENTER)
    add_p(tf, "", 12)
    add_p(tf, "Éditeur français de logiciel de gestion documentaire et de facturation électronique", 14,
          color=GRAY, align=PP_ALIGN.CENTER)
    add_p(tf, "Plan de cession ouvert", 14, color=GRAY, align=PP_ALIGN.CENTER)
    add_p(tf, "", 12)
    add_p(tf, "Brantham Partners, Conseil M&A distressed", 14, align=PP_ALIGN.CENTER, before=12)
    add_p(tf, "2026-05-20", 14, align=PP_ALIGN.CENTER)

    tb(s, 2.0, 6.0, 9.3, 0.6,
       "Document confidentiel, communiqué au repreneur dans le cadre de notre mandat de représentation",
       size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — TEASER (1-page investment summary, B&W booktabs)
# ═══════════════════════════════════════════════════════════════
def teaser(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)

    # Title
    tb(s, 0.8, 0.35, 11.7, 0.6, "1. L'opportunité en bref", size=22, bold=True)
    rule(s, 1.1, l=0.8)

    # KPI table — booktabs (top row = header style but in body)
    tbl_w = Inches(11.7)
    tbl_h = Inches(0.35 * 10)
    ts = s.shapes.add_table(10, 2, Inches(0.8), Inches(1.45), tbl_w, tbl_h)
    tbl = ts.table
    tbl.columns[0].width = int(tbl_w * 0.65)
    tbl.columns[1].width = int(tbl_w * 0.35)

    kpis = [
        ("Indicateur", "Ordre de grandeur"),
        ("Chiffre d'affaires (facturation brute)", "de l'ordre de 9 M€ *"),
        ("Revenu récurrent (abonnements et licences cloud)", "de l'ordre de 7 à 7,5 M€, soit ~ 85 % du CA"),
        ("Prestations (déploiement, maintenance, formation)", "~ 15 % du CA"),
        ("Base installée", "~ 3 000 clients actifs, 250 000+ utilisateurs"),
        ("Réseau de distribution", "~ 60 partenaires distributeurs"),
        ("Effectif", "équipe d'une cinquantaine de personnes"),
        ("Passif non repris", "plusieurs millions d'euros, laissés à la procédure"),
        ("Actif corporel", "marginal, la valeur est quasi intégralement immatérielle"),
        ("Statut réglementaire", "plateforme agréée facture électronique (agrément 2026)"),
    ]

    for i, (a, b) in enumerate(kpis):
        c0 = tbl.cell(i, 0); c0.text = ""
        pp = c0.text_frame.paragraphs[0]
        pp.text = a; pp.font.name = FONT; pp.font.size = Pt(11)
        pp.font.bold = (i == 0); pp.font.color.rgb = BLACK

        c1 = tbl.cell(i, 1); c1.text = ""
        pp = c1.text_frame.paragraphs[0]
        pp.text = b; pp.font.name = FONT; pp.font.size = Pt(11)
        pp.font.bold = (i == 0); pp.font.color.rgb = BLACK
        pp.alignment = PP_ALIGN.RIGHT

    # booktabs borders
    for j in range(2):
        tc = tbl.cell(0, j)._tc; tcPr = tc.get_or_add_tcPr()
        for tag, wt in [('a:lnT', '19050'), ('a:lnB', '9525')]:
            ln = tcPr.find(qn(tag))
            if ln is None: ln = tcPr.makeelement(qn(tag), {}); tcPr.append(ln)
            ln.set('w', wt)
            sf = ln.find(qn('a:solidFill'))
            if sf is None: sf = ln.makeelement(qn('a:solidFill'), {}); ln.append(sf)
            sc = sf.find(qn('a:srgbClr'))
            if sc is None: sc = sf.makeelement(qn('a:srgbClr'), {'val': '000000'}); sf.append(sc)

    for j in range(2):
        tc = tbl.cell(9, j)._tc; tcPr = tc.get_or_add_tcPr()
        ln = tcPr.find(qn('a:lnB'))
        if ln is None: ln = tcPr.makeelement(qn('a:lnB'), {}); tcPr.append(ln)
        ln.set('w', '19050')
        sf = ln.find(qn('a:solidFill'))
        if sf is None: sf = ln.makeelement(qn('a:solidFill'), {}); ln.append(sf)
        sc = sf.find(qn('a:srgbClr'))
        if sc is None: sc = sf.makeelement(qn('a:srgbClr'), {'val': '000000'}); sf.append(sc)

    # Footnote
    tb(s, 0.8, 5.1, 11.7, 0.5,
       "* Les comptes publiés font l'objet de retraitements. Notre analyse reconstruit le CA et la rentabilité réels.",
       size=9, italic=True, color=GRAY)

    # Bottom text
    tb(s, 0.8, 5.35, 11.7, 2.0,
       "La restructuration sociale est déjà engagée par la procédure. "
       "Le passif n'est pas repris : il reste à la charge de la procédure. "
       "Le repreneur acquiert donc une base de revenus récurrents et un actif réglementaire, "
       "sans le poids financier qui a conduit la société en difficulté.",
       size=12, color=BLACK)

    rule(s, 6.8, l=0.8)
    # footer
    tb(s, 0.8, 6.9, 11.7, 0.4,
       "Confidentiel — Brantham Partners — paul.roulleau@branthampartners.fr",
       size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — THÈSE / ACTIFS
# ═══════════════════════════════════════════════════════════════
def these(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)

    tb(s, 0.8, 0.35, 11.7, 0.6, "2. La thèse de reprise : où est la valeur", size=22, bold=True)
    rule(s, 1.1, l=0.8)

    tb(s, 0.8, 1.4, 11.7, 1.0,
       "La cible est en difficulté financière, mais cette difficulté est de nature financière et organisationnelle, "
       "non liée au produit ou au marché. Le plan de cession permet d'acquérir la valeur stratégique en payant "
       "un prix d'actif, le passif étant laissé à la procédure. C'est le ressort central du dossier.",
       size=12)

    # Actifs table
    tbl_w = Inches(11.7)
    tbl_h = Inches(0.4 * 5)
    ts = s.shapes.add_table(5, 2, Inches(0.8), Inches(2.5), tbl_w, tbl_h)
    tbl = ts.table
    tbl.columns[0].width = int(tbl_w * 0.4)
    tbl.columns[1].width = int(tbl_w * 0.6)

    actifs = [
        ("Actif", "Pourquoi il a de la valeur"),
        ("Agrément de facturation électronique", "Droit d'entrée rare sur un marché réglementaire, barrière de plusieurs mois"),
        ("Base installée ~ 3 000 clients", "Revenu récurrent (abonnements cloud), socle prévisible"),
        ("Réseau ~ 60 partenaires", "Canal de distribution indirect, impossible à reconstituer rapidement"),
        ("Plateforme technique et brique IA", "Suite logicielle mature, cloud/on-premise, filiale IA documentaire"),
    ]
    for i, (a, b) in enumerate(actifs):
        c0 = tbl.cell(i, 0); c0.text = ""
        pp = c0.text_frame.paragraphs[0]; pp.text = a
        pp.font.name = FONT; pp.font.size = Pt(11); pp.font.bold = (i == 0); pp.font.color.rgb = BLACK

        c1 = tbl.cell(i, 1); c1.text = ""
        pp = c1.text_frame.paragraphs[0]; pp.text = b
        pp.font.name = FONT; pp.font.size = Pt(11); pp.font.bold = (i == 0); pp.font.color.rgb = BLACK

    # booktabs
    for j in range(2):
        tc = tbl.cell(0, j)._tc; tcPr = tc.get_or_add_tcPr()
        for tag, wt in [('a:lnT', '19050'), ('a:lnB', '9525')]:
            ln = tcPr.find(qn(tag)) or tcPr.makeelement(qn(tag), {})
            tcPr.append(ln) if tcPr.find(qn(tag)) is None else None
            ln.set('w', wt)
            sf = ln.find(qn('a:solidFill')) or ln.makeelement(qn('a:solidFill'), {})
            if ln.find(qn('a:solidFill')) is None: ln.append(sf)
            sc = sf.find(qn('a:srgbClr')) or sf.makeelement(qn('a:srgbClr'), {'val': '000000'})
            if sf.find(qn('a:srgbClr')) is None: sf.append(sc)

    for j in range(2):
        tc = tbl.cell(4, j)._tc; tcPr = tc.get_or_add_tcPr()
        ln = tcPr.find(qn('a:lnB')) or tcPr.makeelement(qn('a:lnB'), {})
        tcPr.append(ln) if tcPr.find(qn('a:lnB')) is None else None
        ln.set('w', '19050')
        sf = ln.find(qn('a:solidFill')) or ln.makeelement(qn('a:solidFill'), {})
        if ln.find(qn('a:solidFill')) is None: ln.append(sf)
        sc = sf.find(qn('a:srgbClr')) or sf.makeelement(qn('a:srgbClr'), {'val': '000000'})
        if sf.find(qn('a:srgbClr')) is None: sf.append(sc)

    # Bottom text
    tb(s, 0.8, 4.7, 11.7, 1.0,
       "Le repreneur n'achète ni le passif (apuré par la procédure) ni la structure de coûts existante. "
       "Il reprend les actifs et les contrats utiles, et reconstruit dès le premier jour une organisation "
       "calibrée sur son projet. C'est une liberté que n'offre pas un rachat de titres classique.",
       size=12)

    rule(s, 6.8, l=0.8)
    tb(s, 0.8, 6.9, 11.7, 0.4,
       "Confidentiel — Brantham Partners — paul.roulleau@branthampartners.fr",
       size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — LEVIERS (booktabs)
# ═══════════════════════════════════════════════════════════════
def leviers(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)

    tb(s, 0.8, 0.35, 11.7, 0.6, "3. Leviers de création de valeur", size=22, bold=True)
    rule(s, 1.1, l=0.8)

    tbl_w = Inches(11.7)
    tbl_h = Inches(0.4 * 7)
    ts = s.shapes.add_table(7, 2, Inches(0.8), Inches(1.5), tbl_w, tbl_h)
    tbl = ts.table
    tbl.columns[0].width = int(tbl_w * 0.35)
    tbl.columns[1].width = int(tbl_w * 0.65)

    rows = [
        ("Levier", "Mécanique"),
        ("Recalibrage du périmètre", "Reprise d'une structure alignée sur le volume d'activité réel"),
        ("Assainissement du bilan", "Disparition du passif et des charges financières"),
        ("Monétisation de l'agrément", "Conversion de la base installée en clients e-facture obligatoire"),
        ("Vente croisée", "Activation de la base sur modules et services à plus forte marge"),
        ("Réactivation du canal partenaires", "Redéploiement commercial via le réseau de distributeurs"),
        ("Synergies d'un acquéreur industriel", "Mutualisation des fonctions support et croisement des portefeuilles"),
    ]
    for i, (a, b) in enumerate(rows):
        c0 = tbl.cell(i, 0); c0.text = ""
        pp = c0.text_frame.paragraphs[0]; pp.text = a
        pp.font.name = FONT; pp.font.size = Pt(11); pp.font.bold = (i == 0); pp.font.color.rgb = BLACK

        c1 = tbl.cell(i, 1); c1.text = ""
        pp = c1.text_frame.paragraphs[0]; pp.text = b
        pp.font.name = FONT; pp.font.size = Pt(11); pp.font.bold = (i == 0); pp.font.color.rgb = BLACK

    # booktabs
    for j in range(2):
        tc = tbl.cell(0, j)._tc; tcPr = tc.get_or_add_tcPr()
        for tag, wt in [('a:lnT', '19050'), ('a:lnB', '9525')]:
            ln = tcPr.find(qn(tag)) or tcPr.makeelement(qn(tag), {})
            if tcPr.find(qn(tag)) is None: tcPr.append(ln)
            ln.set('w', wt)
            sf = ln.find(qn('a:solidFill')) or ln.makeelement(qn('a:solidFill'), {})
            if ln.find(qn('a:solidFill')) is None: ln.append(sf)
            sc = sf.find(qn('a:srgbClr'))
            if sc is None: sc = sf.makeelement(qn('a:srgbClr'), {'val': '000000'}); sf.append(sc)

    for j in range(2):
        tc = tbl.cell(6, j)._tc; tcPr = tc.get_or_add_tcPr()
        ln = tcPr.find(qn('a:lnB')) or tcPr.makeelement(qn('a:lnB'), {})
        if tcPr.find(qn('a:lnB')) is None: tcPr.append(ln)
        ln.set('w', '19050')
        sf = ln.find(qn('a:solidFill')) or ln.makeelement(qn('a:solidFill'), {})
        if ln.find(qn('a:solidFill')) is None: ln.append(sf)
        sc = sf.find(qn('a:srgbClr'))
        if sc is None: sc = sf.makeelement(qn('a:srgbClr'), {'val': '000000'}); sf.append(sc)

    rule(s, 6.8, l=0.8)
    tb(s, 0.8, 6.9, 11.7, 0.4,
       "Confidentiel — Brantham Partners — paul.roulleau@branthampartners.fr",
       size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — SCÉNARIOS
# ═══════════════════════════════════════════════════════════════
def scenarios(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)

    tb(s, 0.8, 0.35, 11.7, 0.6, "4. Scénarios de reprise", size=22, bold=True)
    rule(s, 1.1, l=0.8)

    tb(s, 0.8, 1.4, 11.7, 0.6,
       "Trois trajectoires-types. Les chiffrages détaillés figurent dans notre modèle, restitué au repreneur.",
       size=11, color=GRAY, italic=True)

    # Scenario A
    tb(s, 0.8, 2.2, 11.7, 0.35, "Scénario A — Reprise autonome et redressement opérationnel", size=14, bold=True)
    tb(s, 1.0, 2.65, 11.5, 1.0,
       "Pour un repreneur indépendant. La société est reprise sur un périmètre resserré, recentrée sur sa base "
       "de clients récurrents et son agrément. La rentabilité se reconstruit par la discipline de coûts et "
       "l'assainissement du bilan. La croissance vient ensuite, portée par la facture électronique.",
       size=12)

    # Scenario B
    tb(s, 0.8, 3.8, 11.7, 0.35, "Scénario B — Reprise industrielle par un acteur du secteur", size=14, bold=True)
    tb(s, 1.0, 4.25, 11.5, 1.0,
       "Pour un acteur déjà établi sur la dématérialisation ou la facture électronique. La cible est un accélérateur. "
       "La base de 3 000 clients est immédiatement adressable. Le back-office et l'infrastructure se mutualisent. "
       "Les synergies rapprochent le point d'équilibre dès l'intégration.",
       size=12)

    # Scenario C
    tb(s, 0.8, 5.4, 11.7, 0.35, "Scénario C — Reprise plateforme et consolidation", size=14, bold=True)
    tb(s, 1.0, 5.85, 11.5, 0.8,
       "Reprise élargie aux activités voisines de dématérialisation et d'IA documentaire. Constitution d'une "
       "plateforme intégrée couvrant GED, dématérialisation, facture électronique et IA. Effet d'échelle immédiat.",
       size=12)

    rule(s, 6.8, l=0.8)
    tb(s, 0.8, 6.9, 11.7, 0.4,
       "Confidentiel — Brantham Partners — paul.roulleau@branthampartners.fr",
       size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — POINTS DE VIGILANCE
# ═══════════════════════════════════════════════════════════════
def vigilance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)

    tb(s, 0.8, 0.35, 11.7, 0.6, "5. Points de vigilance identifiés", size=22, bold=True)
    rule(s, 1.1, l=0.8)

    items = [
        ("Propriété intellectuelle.",
         "La structure de détention du logiciel et de la marque appelle une vérification précise "
         "et une sécurisation préalable au dépôt de l'offre."),
        ("Contrats critiques.",
         "Certains contrats sont indispensables à la continuité d'exploitation et doivent être "
         "sécurisés dans le cadre de l'offre."),
        ("Périmètre social et calendrier.",
         "La définition du périmètre repris et la tenue du calendrier de la procédure sont des "
         "paramètres clés de la réussite."),
    ]
    y = 1.5
    for label, desc in items:
        tf = tb(s, 0.8, y, 11.7, 0.8, label, size=13, bold=True)
        add_p(tf, desc, size=12, color=BLACK, before=2)
        y += 1.1

    tb(s, 0.8, 4.9, 11.7, 1.0,
       "Identifier ces points, les chiffrer et les sécuriser est ce qui sépare une offre retenue d'une offre "
       "rejetée, et une reprise solide d'une reprise qui se fragilise après le closing.",
       size=12)

    rule(s, 6.8, l=0.8)
    tb(s, 0.8, 6.9, 11.7, 0.4,
       "Confidentiel — Brantham Partners — paul.roulleau@branthampartners.fr",
       size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — BRANTHAM + PROCHAINE ÉTAPE
# ═══════════════════════════════════════════════════════════════
def brantham(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)

    tb(s, 0.8, 0.35, 11.7, 0.6, "6. Ce que Brantham Partners a réalisé et apporte", size=22, bold=True)
    rule(s, 1.1, l=0.8)

    items = [
        "Analyse exhaustive de la data room (> 130 documents : comptable, contractuel, juridique, social).",
        "Reconstruction du modèle économique réel et retraité de la cible, au-delà des comptes publiés.",
        "Modélisation des scénarios de reprise et du périmètre viable, chiffrés.",
        "Identification et plan de sécurisation des points de vigilance.",
    ]
    y = 1.5
    for item in items:
        tb(s, 1.0, y, 11.5, 0.35, "—  " + item, size=12)
        y += 0.45

    tb(s, 0.8, 3.5, 11.7, 1.0,
       "Brantham Partners accompagne ensuite le repreneur de bout en bout : restitution de l'analyse, "
       "construction du périmètre et de l'offre, dépôt auprès des organes de la procédure, accompagnement "
       "à l'audience et suivi post-cession. Notre rémunération est alignée sur la réussite de l'opération.",
       size=12)

    rule(s, 4.9, l=0.8)

    tb(s, 0.8, 5.1, 11.7, 0.35, "Prochaine étape", size=14, bold=True)
    tb(s, 1.0, 5.55, 11.5, 1.0,
       "Restitution détaillée de notre analyse à brève échéance. Le calendrier de la procédure est resserré : "
       "la date limite de dépôt des offres est proche. Un positionnement rapide est la condition d'une offre crédible.",
       size=12)

    rule(s, 6.8, l=0.8)
    tb(s, 0.8, 6.9, 11.7, 0.4,
       "Confidentiel — Brantham Partners — paul.roulleau@branthampartners.fr",
       size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    cover(prs)
    teaser(prs)
    these(prs)
    leviers(prs)
    scenarios(prs)
    vigilance(prs)
    brantham(prs)

    prs.save(OUTPUT)
    print(f"✅ {OUTPUT}")


if __name__ == "__main__":
    build()
