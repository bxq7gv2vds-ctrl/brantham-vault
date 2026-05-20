"""Generate Brantham corporate PowerPoint template (.potx ready).

Style: CMU Serif, noir & blanc, booktabs, épuré — identique à la note stratégique PDF.
Output: vault/08_Exports/Brantham-Corporate-Template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FONT_NAME = "CMU Serif"
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUTPUT = os.path.expanduser(
    "~/vault/08_Exports/Brantham-Corporate-Template.pptx"
)


def set_slide_bg_white(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, size=18,
                bold=False, italic=False, color=BLACK, align=PP_ALIGN.LEFT,
                font_name=FONT_NAME, spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)

    if spacing != 1.0:
        p.line_spacing = Pt(size * spacing)

    return tf


def add_paragraph(tf, text, size=18, bold=False, italic=False,
                  color=BLACK, align=PP_ALIGN.LEFT, spacing=1.15,
                  space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if spacing != 1.0:
        p.line_spacing = Pt(size * spacing)
    return p


def add_page_number(slide, number):
    tf = add_textbox(
        slide, Inches(0.5), Inches(7.0), Inches(2), Inches(0.4),
        str(number), size=10, color=MEDIUM_GRAY,
        align=PP_ALIGN.LEFT
    )


def add_horizontal_rule(slide, left, top, width=Inches(12.333)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.fill.background()
    return shape


def add_cover(slide):
    set_slide_bg_white(slide)
    tf = add_textbox(
        slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1),
        "Titre du document", size=28, bold=True,
        align=PP_ALIGN.CENTER
    )
    add_paragraph(tf, "Sous-titre ou description", size=16, color=DARK_GRAY,
                  align=PP_ALIGN.CENTER, space_before=8)
    add_horizontal_rule(slide, Inches(4), Inches(3.8), Inches(5.3))

    add_textbox(
        slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
        "Brantham Partners — Conseil M&A", size=14, color=DARK_GRAY,
        align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.5),
        "2026-05-20", size=14, color=DARK_GRAY,
        align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.8),
        "Document confidentiel — Ne pas diffuser sans autorisation",
        size=11, italic=True, color=MEDIUM_GRAY,
        align=PP_ALIGN.CENTER
    )


def add_section_divider(slide, number, title):
    set_slide_bg_white(slide)
    add_horizontal_rule(slide, Inches(2), Inches(3.0), Inches(9.3))

    add_textbox(
        slide, Inches(2), Inches(1.8), Inches(9.3), Inches(0.6),
        number, size=16, color=DARK_GRAY, align=PP_ALIGN.LEFT
    )
    add_textbox(
        slide, Inches(2), Inches(3.2), Inches(9.3), Inches(1.2),
        title, size=28, bold=True, align=PP_ALIGN.LEFT
    )

    add_page_number(slide, "")


def add_content_slide(slide, title, body_lines, slide_num):
    set_slide_bg_white(slide)
    add_textbox(
        slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.7),
        title, size=22, bold=True
    )
    add_horizontal_rule(slide, Inches(1), Inches(1.3))

    y = Inches(1.7)
    for line in body_lines:
        if line.startswith("## "):
            add_textbox(
                slide, Inches(1), y, Inches(11.3), Inches(0.4),
                line[3:], size=16, bold=True
            )
            y += Inches(0.5)
        elif line.startswith("### "):
            add_textbox(
                slide, Inches(1.2), y, Inches(11.1), Inches(0.35),
                line[4:], size=14, bold=True, italic=True
            )
            y += Inches(0.45)
        elif line == "":
            y += Inches(0.2)
        else:
            tf = add_textbox(
                slide, Inches(1.2), y, Inches(11.1), Inches(1.2),
                line, size=13, color=DARK_GRAY
            )
            y += Inches(0.45)

    add_page_number(slide, slide_num)


def add_two_column_slide(slide, title, left_items, right_items, slide_num):
    set_slide_bg_white(slide)
    add_textbox(
        slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.7),
        title, size=22, bold=True
    )
    add_horizontal_rule(slide, Inches(1), Inches(1.3))

    tf = add_textbox(
        slide, Inches(1), Inches(1.7), Inches(5.5), Inches(5),
        "", size=13, color=DARK_GRAY
    )
    first = True
    for item in left_items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_NAME
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    tf2 = add_textbox(
        slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(5),
        "", size=13, color=DARK_GRAY
    )
    first = True
    for item in right_items:
        if first:
            p = tf2.paragraphs[0]
            first = False
        else:
            p = tf2.add_paragraph()
        p.text = item
        p.font.name = FONT_NAME
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    add_page_number(slide, slide_num)


def add_booktabs_table(slide, title, headers, rows,
                       col_widths=None, slide_num=""):
    set_slide_bg_white(slide)
    add_textbox(
        slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.7),
        title, size=22, bold=True
    )
    add_horizontal_rule(slide, Inches(1), Inches(1.3))

    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_left = Inches(1)
    tbl_top = Inches(1.7)
    tbl_w = Inches(11.3)
    tbl_h = Inches(0.5 * n_rows)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h
    )
    table = table_shape.table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(tbl_w * w / total)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BLACK
        if h in ("", " "):
            pass

    # Booktabs borders via XML hacking
    from pptx.oxml.ns import qn
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.name = FONT_NAME
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(2)
            p.space_after = Pt(2)

    # Header bottom border (thick)
    for j in range(n_cols):
        tc = table.cell(0, j)._tc
        tcPr = tc.get_or_add_tcPr()
        for border_tag in ['a:lnB']:
            ln = tcPr.find(qn(border_tag))
            if ln is None:
                ln = tcPr.makeelement(qn(border_tag), {})
                tcPr.append(ln)
                ln.set('w', '12700')
                solidFill = ln.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '000000'})
                solidFill.append(srgbClr)
                ln.append(solidFill)

    add_page_number(slide, slide_num)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    # ---- Slide 1: Cover ----
    s1 = prs.slides.add_slide(blank_layout)
    add_cover(s1)

    # ---- Slide 2: Section divider ----
    s2 = prs.slides.add_slide(blank_layout)
    add_section_divider(s2, "01", "Contexte de l'opération")

    # ---- Slide 3: Content ----
    s3 = prs.slides.add_slide(blank_layout)
    add_content_slide(s3,
        "Synthèse de l'opportunité",
        [
            "Un éditeur français de logiciel de gestion documentaire, 22 ans d'existence, ~3 000 clients.",
            "Plan de cession ouvert — reprise par acquisition d'actifs, sans le passif.",
            "## Indicateurs clés",
            "Chiffre d'affaires : ~9 M€, dont ~85 % de revenu récurrent.",
            "Agrément plateforme de facturation électronique obtenu en 2026.",
            "Base installée : 3 000 clients actifs, réseau de 60 partenaires.",
            "## Fenêtre stratégique",
            "La généralisation de la e-facture crée une demande captive.",
            "L'agrément est un droit d'entrée rare sur un marché en formation.",
        ],
        3
    )

    # ---- Slide 4: Two-column ----
    s4 = prs.slides.add_slide(blank_layout)
    add_two_column_slide(s4,
        "Actifs repris vs. Passif laissé",
        [
            "Actifs incorporels",
            "Agrément facture électronique",
            "Base 3 000 clients",
            "Réseau 60 partenaires",
            "Plateforme technique + IA",
            "Contrats de travail repris",
        ],
        [
            "Passif antérieur non repris",
            "Dettes fournisseurs",
            "Dettes fiscales et sociales",
            "Charges financières",
            "Contentieux antérieurs",
            "Passif apuré par la procédure",
        ],
        4
    )

    # ---- Slide 5: Table ----
    s5 = prs.slides.add_slide(blank_layout)
    add_booktabs_table(s5,
        "Scénarios de reprise",
        ["Scénario", "Logique", "Profil"],
        [
            ["A — Autonome", "Redressement opérationnel 18-24 mois", "Entrepreneur, ticket maîtrisé"],
            ["B — Industriel", "Synergies + vente croisée", "Acteur du secteur"],
            ["C — Plateforme", "Consolidation + effet d'échelle", "Fonds ou groupe"],
        ],
        col_widths=[3, 5, 4],
        slide_num=5
    )

    # ---- Slide 6: Content ----
    s6 = prs.slides.add_slide(blank_layout)
    add_content_slide(s6,
        "Prochaines étapes",
        [
            "1. Restitution détaillée de l'analyse avec le repreneur.",
            "2. Transmission des éléments nominatifs du dossier sous NDA.",
            "3. Construction et calibrage de l'offre avec nos équipes.",
            "4. Dépôt de l'offre auprès des organes de la procédure.",
            "5. Accompagnement jusqu'à l'audience et au closing.",
            "",
            "Le calendrier de la procédure est resserré.",
            "Un positionnement rapide est la condition d'une offre crédible.",
        ],
        6
    )

    # ---- Slide 7: Closing / Contact ----
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg_white(s7)
    add_horizontal_rule(s7, Inches(4), Inches(2.5), Inches(5.3))
    add_textbox(
        s7, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.6),
        "Contact", size=24, bold=True, align=PP_ALIGN.CENTER
    )
    tf = add_textbox(
        s7, Inches(1.5), Inches(3.6), Inches(10.3), Inches(2),
        "Paul ROULLEAU", size=18, bold=True,
        align=PP_ALIGN.CENTER
    )
    add_paragraph(tf, "Brantham Partners", size=14, color=DARK_GRAY,
                  align=PP_ALIGN.CENTER, space_before=4)
    add_paragraph(tf, "paul.roulleau@branthampartners.fr", size=14,
                  color=DARK_GRAY, align=PP_ALIGN.CENTER, space_before=2)
    add_paragraph(tf, "07 68 36 25 63", size=14, color=DARK_GRAY,
                  align=PP_ALIGN.CENTER, space_before=2)
    add_textbox(
        s7, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.6),
        "Document confidentiel — Ne pas diffuser sans autorisation",
        size=10, italic=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER
    )

    prs.save(OUTPUT)
    print(f"Template créé : {OUTPUT}")


if __name__ == "__main__":
    build()