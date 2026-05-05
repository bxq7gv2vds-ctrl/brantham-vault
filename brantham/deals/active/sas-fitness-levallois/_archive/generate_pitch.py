"""Generate pitch-repreneur.pptx — Salle de sport Levallois.

DA strictly follows the LinkedIn carousel grammar:
- Cream background (#FAFAF8), ink text (#0F0F0E)
- Inter sans-serif body, Instrument Serif italic for inline emphasis, DM Mono for eyebrows / data
- Strict typographic hierarchy, generous margins, no decorative shapes
- Closer slides 12-13 inverted (ink background, cream text)

Run:
    python3 generate_pitch.py
Output:
    pitch-repreneur.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ------------------------------------------------------------------
# Palette + typography
# ------------------------------------------------------------------
INK = RGBColor(0x0F, 0x0F, 0x0E)
INK3 = RGBColor(0x4A, 0x4A, 0x47)
INK4 = RGBColor(0x8A, 0x8A, 0x86)
INK5 = RGBColor(0xB8, 0xB8, 0xB4)
CREAM = RGBColor(0xFA, 0xFA, 0xF8)
OFF = RGBColor(0xF4, 0xF3, 0xF0)
INDIGO = RGBColor(0x5E, 0x54, 0xF0)
RED = RGBColor(0xC8, 0x25, 0x1A)
GREEN = RGBColor(0x1A, 0x7A, 0x3A)

SANS = "Inter"
SERIF = "Instrument Serif"
MONO = "DM Mono"

# ------------------------------------------------------------------
# Geometry — 16:9, 13.333" × 7.5"
# ------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LEFT = Inches(0.85)        # left margin
TOP = Inches(0.55)         # eyebrow vertical
RIGHT = Inches(0.85)
CONTENT_TOP = Inches(1.4)  # where titles start


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def style_run(run, *, font=SANS, size=14, bold=False, italic=False, color=INK,
              tracking=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # Tracking via XML (python-pptx doesn't expose it directly)
    if tracking is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(tracking)))


def add_eyebrow(slide, text, *, dark=False):
    color_dot = INDIGO if not dark else RGBColor(0xB0, 0xB0, 0xAE)
    color_txt = INK4 if not dark else RGBColor(0x9A, 0x9A, 0x97)
    # dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, LEFT, TOP + Emu(50000),
                                 Inches(0.09), Inches(0.09))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color_dot
    dot.line.fill.background()
    # text
    tb, tf = add_textbox(slide, LEFT + Inches(0.22), TOP - Emu(10000),
                         Inches(8), Inches(0.3))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    style_run(run, font=MONO, size=10, color=color_txt, tracking=200)


def add_pagination(slide, num, total, *, dark=False):
    color = INK5 if not dark else RGBColor(0x70, 0x70, 0x6E)
    tb, tf = add_textbox(slide, SLIDE_W - RIGHT - Inches(2),
                         TOP - Emu(10000), Inches(2), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num:02d} / {total:02d}"
    style_run(run, font=MONO, size=10, color=color, tracking=200)


def add_footer(slide, label, *, dark=False, url=None):
    line_color = RGBColor(0xE5, 0xE4, 0xE0) if not dark else RGBColor(0x2A, 0x2A, 0x28)
    line = slide.shapes.add_connector(1, LEFT, SLIDE_H - Inches(0.65),
                                       SLIDE_W - RIGHT, SLIDE_H - Inches(0.65))
    line.line.color.rgb = line_color
    line.line.width = Pt(0.5)

    # logo mark (small black square with "A")
    mark_size = Inches(0.32)
    mark_color = INK if not dark else CREAM
    text_color_in_mark = CREAM if not dark else INK
    mark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT,
                                  SLIDE_H - Inches(0.5), mark_size, mark_size)
    mark.fill.solid()
    mark.fill.fore_color.rgb = mark_color
    mark.line.fill.background()
    mtf = mark.text_frame
    mtf.margin_left = mtf.margin_right = Emu(0)
    mtf.margin_top = mtf.margin_bottom = Emu(0)
    mtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    mp = mtf.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "A"
    style_run(mr, font=SANS, size=14, bold=True, color=text_color_in_mark)

    # name
    name_color = INK if not dark else CREAM
    tb, tf = add_textbox(slide, LEFT + Inches(0.45),
                         SLIDE_H - Inches(0.48), Inches(4), Inches(0.3))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Brantham Partners"
    style_run(run, font=SANS, size=12, color=name_color)

    # right label
    right_color = INK5 if not dark else RGBColor(0x70, 0x70, 0x6E)
    tb, tf = add_textbox(slide, SLIDE_W - RIGHT - Inches(5),
                         SLIDE_H - Inches(0.46), Inches(5), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = (url or label).upper()
    style_run(run, font=MONO, size=9, color=right_color, tracking=200)


def add_accent_line(slide, top, *, color=INK, width=Inches(0.6)):
    line = slide.shapes.add_connector(1, LEFT, top, LEFT + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)


def add_title(slide, top, segments, size=42, line_height=1.05):
    """segments = list of (text, italic_serif: bool, line_break_after: bool)."""
    tb, tf = add_textbox(slide, LEFT, top,
                         SLIDE_W - LEFT - RIGHT, Inches(2.5))
    p = tf.paragraphs[0]
    p.line_spacing = line_height
    for i, (text, serif, br) in enumerate(segments):
        run = p.add_run()
        run.text = text
        if serif:
            style_run(run, font=SERIF, size=size, italic=True, color=INK)
        else:
            style_run(run, font=SANS, size=size, color=INK)
        if br:
            p = tf.add_paragraph()
            p.line_spacing = line_height


def add_subtitle(slide, top, text, *, color=INK3, size=14, max_width=Inches(8)):
    tb, tf = add_textbox(slide, LEFT, top, max_width, Inches(2))
    p = tf.paragraphs[0]
    p.line_spacing = 1.55
    run = p.add_run()
    run.text = text
    style_run(run, font=SANS, size=size, color=color)


def add_step(slide, top, num, title, desc, *, qualifier=None, qualifier_color=INK4,
             amount=None, amount_color=INK, dark=False):
    """One numbered list row, LinkedIn-style.
    Returns the bottom y (Emu) so subsequent rows can stack.
    """
    row_h = Inches(1.0)
    label_color = INK if not dark else CREAM
    desc_color = INK3 if not dark else RGBColor(0xB0, 0xB0, 0xAE)
    num_color = INK4 if not dark else RGBColor(0x80, 0x80, 0x7E)
    line_color = RGBColor(0xE0, 0xDF, 0xDB) if not dark else RGBColor(0x2A, 0x2A, 0x28)

    # number
    tb, tf = add_textbox(slide, LEFT, top + Inches(0.06),
                         Inches(0.5), Inches(0.4))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{num:02d}"
    style_run(run, font=MONO, size=11, color=num_color, tracking=180)

    # title + desc
    text_left = LEFT + Inches(0.6)
    text_w = SLIDE_W - text_left - RIGHT - Inches(2.6)
    tb, tf = add_textbox(slide, text_left, top, text_w, row_h)
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    run = p.add_run()
    run.text = title
    style_run(run, font=SANS, size=18, bold=False, color=label_color)
    run.font.bold = True
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.4
    p2.space_before = Pt(4)
    run = p2.add_run()
    run.text = desc
    style_run(run, font=SANS, size=13, color=desc_color)

    # right side qualifier OR amount
    right_w = Inches(2.4)
    right_left = SLIDE_W - RIGHT - right_w
    tb, tf = add_textbox(slide, right_left, top + Inches(0.08),
                         right_w, Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    if amount is not None:
        run = p.add_run()
        run.text = amount
        style_run(run, font=MONO, size=18, bold=True, color=amount_color)
    elif qualifier is not None:
        run = p.add_run()
        run.text = qualifier.upper()
        style_run(run, font=MONO, size=10, color=qualifier_color, tracking=160)

    # bottom line
    line_y = top + row_h + Inches(0.05)
    line = slide.shapes.add_connector(1, LEFT, line_y,
                                       SLIDE_W - RIGHT, line_y)
    line.line.color.rgb = line_color
    line.line.width = Pt(0.5)

    return line_y + Inches(0.2)


def add_alert(slide, top, label, text, *, color=RED, max_w=Inches(10.5)):
    """Pastel encart à la 'À surveiller'."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT, top,
                                 max_w, Inches(1.2))
    box.adjustments[0] = 0.08
    box.fill.solid()
    fill_rgb = RGBColor(min(color[0]+50, 255), min(color[1]+50, 255), min(color[2]+50, 255)) \
        if False else color
    # use a soft tint manually
    if color == RED:
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xEC, 0xEA)
    elif color == GREEN:
        box.fill.fore_color.rgb = RGBColor(0xEC, 0xF6, 0xEF)
    else:
        box.fill.fore_color.rgb = OFF
    box.line.color.rgb = RGBColor(0xE5, 0xE4, 0xE0) if color == INK3 else \
        RGBColor(0xE8, 0xCE, 0xCB) if color == RED else \
        RGBColor(0xC9, 0xE1, 0xD1)
    box.line.width = Pt(0.5)

    tb, tf = add_textbox(slide, LEFT + Inches(0.35), top + Inches(0.2),
                         max_w - Inches(0.7), Inches(1))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label.upper()
    style_run(run, font=MONO, size=10, color=color if color != INK3 else INK4,
              tracking=200)
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.4
    p2.space_before = Pt(8)
    run = p2.add_run()
    run.text = text
    style_run(run, font=SANS, size=15, color=INK)
    return top + Inches(1.4)


def add_tag(slide, left, top, text, *, color=INK3, border=RGBColor(0xE5, 0xE4, 0xE0),
            bg=OFF):
    """Returns the right edge x of the tag for chaining."""
    # Approximate width based on text length
    chars = len(text)
    w = Inches(0.35) + Inches(0.085 * chars)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w,
                                  Inches(0.36))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg
    pill.line.color.rgb = border
    pill.line.width = Pt(0.5)
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    style_run(run, font=MONO, size=9, color=color, tracking=140)
    return left + w + Inches(0.12)


def add_cascade_row(slide, top, num, label, amount_eur, pct, *, red=False):
    """Money cascade row with horizontal bar."""
    row_h = Inches(0.85)
    # number
    tb, tf = add_textbox(slide, LEFT, top + Inches(0.1), Inches(0.5), Inches(0.4))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{num:02d}"
    style_run(run, font=MONO, size=11, color=INK4, tracking=180)
    # label
    tb, tf = add_textbox(slide, LEFT + Inches(0.6), top, Inches(7), Inches(0.4))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    style_run(run, font=SANS, size=17, bold=True, color=INK)
    # bar
    bar_max_w = Inches(7)
    bar_w = Inches(7 * pct / 100)
    bar_y = top + Inches(0.5)
    bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT + Inches(0.6),
                                    bar_y, bar_max_w, Inches(0.18))
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = RGBColor(0xEE, 0xED, 0xE9)
    bg_bar.line.fill.background()
    fg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT + Inches(0.6),
                                    bar_y, bar_w if bar_w > 0 else Emu(1),
                                    Inches(0.18))
    fg_bar.fill.solid()
    fg_bar.fill.fore_color.rgb = RED if red else INK
    fg_bar.line.fill.background()
    # amount
    right_w = Inches(2)
    right_left = SLIDE_W - RIGHT - right_w
    tb, tf = add_textbox(slide, right_left, top + Inches(0.1), right_w, Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = amount_eur
    style_run(run, font=MONO, size=24, bold=True, color=INK)
    # bottom separator
    line_y = top + row_h + Inches(0.05)
    line = slide.shapes.add_connector(1, LEFT, line_y, SLIDE_W - RIGHT, line_y)
    line.line.color.rgb = RGBColor(0xE5, 0xE4, 0xE0)
    line.line.width = Pt(0.5)
    return line_y + Inches(0.15)


def add_data_table(slide, top, headers, rows):
    """Simple data table — header row + data rows, all DM Mono.
    rows = list of (label, [(value, color), ...])
    """
    n_cols = len(headers)
    col_w = (SLIDE_W - LEFT - RIGHT) / n_cols

    # headers
    for i, h in enumerate(headers):
        tb, tf = add_textbox(slide, LEFT + col_w * i, top, col_w, Inches(0.4))
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = h.upper()
        style_run(run, font=MONO, size=10, color=INK4, tracking=180)

    line_y = top + Inches(0.5)
    line = slide.shapes.add_connector(1, LEFT, line_y, SLIDE_W - RIGHT, line_y)
    line.line.color.rgb = RGBColor(0xD8, 0xD7, 0xD3)
    line.line.width = Pt(0.6)

    y = line_y + Inches(0.15)
    for label, values in rows:
        tb, tf = add_textbox(slide, LEFT, y, col_w, Inches(0.5))
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        style_run(run, font=SANS, size=15, color=INK3)
        for i, (val, color) in enumerate(values):
            tb, tf = add_textbox(slide, LEFT + col_w * (i + 1), y, col_w,
                                 Inches(0.5))
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            style_run(run, font=MONO, size=18, bold=True, color=color)
        y += Inches(0.55)
        line = slide.shapes.add_connector(1, LEFT, y, SLIDE_W - RIGHT, y)
        line.line.color.rgb = RGBColor(0xE5, 0xE4, 0xE0)
        line.line.width = Pt(0.4)
        y += Inches(0.15)
    return y


def add_two_col(slide, top, left_label, left_rows, right_label, right_rows,
                *, right_accent=GREEN):
    """Two columns side-by-side, LinkedIn signal-style."""
    gap = Inches(0.4)
    col_w = (SLIDE_W - LEFT - RIGHT - gap) / 2
    col_h = Inches(3.4)

    def render_col(x, label, rows, accent=None):
        bg = OFF if accent is None else (
            RGBColor(0xEC, 0xF6, 0xEF) if accent == GREEN else RGBColor(0xFA, 0xEC, 0xEA)
        )
        border = RGBColor(0xE5, 0xE4, 0xE0) if accent is None else (
            RGBColor(0xC9, 0xE1, 0xD1) if accent == GREEN else RGBColor(0xE8, 0xCE, 0xCB)
        )
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top,
                                     col_w, col_h)
        box.adjustments[0] = 0.04
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = border
        box.line.width = Pt(0.5)

        tb, tf = add_textbox(slide, x + Inches(0.4), top + Inches(0.3),
                             col_w - Inches(0.8), Inches(0.4))
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label.upper()
        label_color = INK4 if accent is None else (GREEN if accent == GREEN else RED)
        style_run(run, font=MONO, size=10, color=label_color, tracking=200)

        ry = top + Inches(0.85)
        for k, v, vcolor in rows:
            tb, tf = add_textbox(slide, x + Inches(0.4), ry,
                                 col_w - Inches(0.8), Inches(0.4))
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = k
            style_run(run, font=SANS, size=13, color=INK)

            tb, tf = add_textbox(slide, x + Inches(0.4), ry,
                                 col_w - Inches(0.8), Inches(0.4))
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = v
            style_run(run, font=MONO, size=13, bold=True,
                      color=vcolor if vcolor else INK)
            ry += Inches(0.42)
            ln = slide.shapes.add_connector(1, x + Inches(0.4), ry,
                                             x + col_w - Inches(0.4), ry)
            ln.line.color.rgb = RGBColor(0xE0, 0xDF, 0xDB)
            ln.line.width = Pt(0.4)
            ry += Inches(0.05)

    render_col(LEFT, left_label, left_rows, accent=None if right_accent == GREEN else None)
    render_col(LEFT + col_w + gap, right_label, right_rows, accent=right_accent)


# ------------------------------------------------------------------
# Build presentation
# ------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ---------- 01 Cover ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Présentation repreneur")
    add_pagination(s, 1, 13)
    add_title(s, Inches(2.6), [
        ("Reprendre une ", False, False),
        ("salle de sport", True, True),
        ("en redressement judiciaire.", False, False),
    ], size=58)
    add_subtitle(s, Inches(4.5),
                 "SAS Fitness Levallois — dossier d'analyse, 28 avril 2026.\n"
                 "Préparé par Brantham Partners pour le repreneur.",
                 size=15, max_width=Inches(9))
    # tags
    x = LEFT
    x = add_tag(s, x, Inches(5.7), "Levallois 92")
    x = add_tag(s, x, Inches(5.7), "CA 350 K€")
    x = add_tag(s, x, Inches(5.7), "DLDO 21/05/2026  12h",
                color=RED, border=RGBColor(0xE8, 0xCE, 0xCB),
                bg=RGBColor(0xFA, 0xEC, 0xEA))
    add_footer(s, "Confidentiel — usage interne")

    # ---------- 02 Identité ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Cible")
    add_pagination(s, 2, 13)
    add_title(s, CONTENT_TOP, [
        ("La société, ", False, False),
        ("en deux minutes", True, False),
        (".", False, False),
    ], size=44)
    add_accent_line(s, Inches(2.65))
    y = Inches(3.0)
    y = add_step(s, y, 1, "Activité",
                 "Salle de sport indépendante — 900 m² — Levallois-Perret (92)",
                 qualifier="Coaching · cardio · musculation")
    y = add_step(s, y, 2, "Forme juridique",
                 "SAS Fitness Levallois — enseigne O3 Partners",
                 qualifier="Dirigeant unique")
    y = add_step(s, y, 3, "Effectif",
                 "[N] salariés — masse salariale annuelle [Y] K€",
                 qualifier="Bulletins de paie en DR")
    y = add_step(s, y, 4, "Chiffre d'affaires 2024",
                 "350 K€ vs pic 2022 à 595 K€ — perte 41 % en deux exercices",
                 qualifier="Source liasse fiscale")
    add_footer(s, "Acte 1 · Contexte")

    # ---------- 03 Procédure ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Procédure")
    add_pagination(s, 3, 13)
    add_title(s, CONTENT_TOP, [
        ("Redressement judiciaire ouvert", False, True),
        ("le ", False, False),
        ("[date d'ouverture]", True, False),
        (".", False, False),
    ], size=44)
    add_accent_line(s, Inches(3.0))
    y = Inches(3.3)
    y = add_step(s, y, 1, "Tribunal",
                 "Tribunal de commerce de Nanterre",
                 qualifier="Hauts-de-Seine")
    y = add_step(s, y, 2, "Administrateur judiciaire",
                 "[Maître X] — cabinet [Y]",
                 qualifier="Mission cession")
    y = add_step(s, y, 3, "Mandataire judiciaire",
                 "[Maître Z]",
                 qualifier="Représentant créanciers")
    y = add_step(s, y, 4, "Échéance dépôt offres",
                 "Greffe Nanterre — pli cacheté avec garantie bancaire",
                 amount="21/05/2026 · 12h", amount_color=RED)
    add_footer(s, "Acte 1 · Contexte")

    # ---------- 04 Trajectoire ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Trajectoire")
    add_pagination(s, 4, 13)
    add_title(s, CONTENT_TOP, [
        ("Quatre ans, ", False, False),
        ("une descente continue", True, False),
        (".", False, False),
    ], size=44)
    add_accent_line(s, Inches(2.65))
    add_data_table(s, Inches(3.05),
                   ["Indicateur", "2022", "2023", "2024", "PO 6 mois"],
                   [
                       ("Chiffre d'affaires",
                        [("595 K€", INK), ("478 K€", INK), ("350 K€", INK), ("175 K€", INK)]),
                       ("EBE",
                        [("62 K€", INK), ("14 K€", INK), ("−38 K€", RED), ("+12 K€", GREEN)]),
                       ("Résultat net",
                        [("21 K€", INK), ("−24 K€", RED), ("−72 K€", RED), ("+4 K€", GREEN)]),
                       ("Trésorerie fin de période",
                        [("48 K€", INK), ("22 K€", INK), ("−9 K€", RED), ("+18 K€", GREEN)]),
                   ])
    # footnote
    tb, tf = add_textbox(s, LEFT, Inches(6.4), Inches(10), Inches(0.3))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Chiffres normalisés — source liasses fiscales + note expert-comptable".upper()
    style_run(run, font=MONO, size=9, color=INK5, tracking=180)
    add_footer(s, "Acte 2 · Diagnostic")

    # ---------- 05 Point vital ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Point vital")
    add_pagination(s, 5, 13)
    # giant red 01
    tb, tf = add_textbox(s, LEFT, Inches(1.2), Inches(3), Inches(2))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "01"
    style_run(run, font=SERIF, size=140, italic=False, color=RED)
    add_title(s, Inches(3.0), [
        ("Le bail commercial", False, True),
        ("a expiré le 30/06/2024", True, False),
        (".", False, False),
    ], size=42)
    add_accent_line(s, Inches(4.6), color=RED)
    add_subtitle(s, Inches(4.85),
                 "Aucune reconduction tacite signée. Le local est occupé sans titre depuis 22 mois.",
                 size=14, max_width=Inches(11))
    add_alert(s, Inches(5.5), "À savoir",
              "La bailleresse est elle-même demanderesse au redressement judiciaire. Sans renégociation d'un nouveau bail à conditions soutenables, aucune reprise n'est viable. C'est l'unique levier binaire du dossier.",
              color=RED, max_w=Inches(11.5))
    add_footer(s, "Acte 2 · Diagnostic")

    # ---------- 06 PO récente ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Période d'observation")
    add_pagination(s, 6, 13)
    add_title(s, CONTENT_TOP, [
        ("6 mois sous AJ,", False, True),
        ("le cash redevient ", False, False),
        ("positif", True, False),
        (".", False, False),
    ], size=44)
    add_accent_line(s, Inches(3.0))
    y = Inches(3.3)
    y = add_step(s, y, 1, "Cash net 6 mois normalisé",
                 "Avant prise en compte du loyer plein — après gel des dettes antérieures",
                 amount="+12 K€", amount_color=GREEN)
    y = add_step(s, y, 2, "Annualisé",
                 "Extrapolation linéaire sur 12 mois, conditions PO inchangées",
                 amount="+28 K€", amount_color=GREEN)
    y = add_step(s, y, 3, "Si bail renégocié à 110 K€/an",
                 "Hypothèse cible — voir slide 08",
                 amount="+98 K€", amount_color=GREEN)
    add_alert(s, y + Inches(0.1), "Lecture",
              "L'exploitation est opérationnellement saine. Ce qui a tué la société, c'est le loyer, pas le business.",
              color=INK3, max_w=Inches(11.5))
    add_footer(s, "Acte 2 · Diagnostic")

    # ---------- 07 Money cascade ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Flux")
    add_pagination(s, 7, 13)
    add_title(s, CONTENT_TOP, [
        ("Sur 100 € encaissés,", False, True),
        ("où part l'argent", True, False),
        (".", False, False),
    ], size=44)
    add_accent_line(s, Inches(3.0))
    y = Inches(3.3)
    y = add_cascade_row(s, y, 1, "Loyer commercial", "51 €", 51, red=True)
    y = add_cascade_row(s, y, 2, "Masse salariale chargée", "28 €", 28)
    y = add_cascade_row(s, y, 3, "Charges fixes (énergie, assurance, télécom)", "14 €", 14)
    y = add_cascade_row(s, y, 4, "Achats & fluides", "5 €", 5)
    # footnote
    tb, tf = add_textbox(s, LEFT, Inches(6.4), Inches(10), Inches(0.3))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Base CA 2024 — structure avant renégociation du bail".upper()
    style_run(run, font=MONO, size=9, color=INK5, tracking=180)
    add_footer(s, "Acte 3 · Argent")

    # ---------- 08 Levier ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Levier")
    add_pagination(s, 8, 13)
    add_title(s, CONTENT_TOP, [
        ("Renégocier le loyer,", False, True),
        ("seule variable qui change tout", True, False),
        (".", False, False),
    ], size=42)
    add_accent_line(s, Inches(3.0))
    add_two_col(s, Inches(3.3),
                "Avant — bail historique",
                [("Loyer annuel", "180 K€", INK),
                 ("% du CA", "51 %", INK),
                 ("EBE An 1", "−12 K€", RED),
                 ("Cash net An 1", "−8 K€", RED)],
                "Après — bail renégocié",
                [("Loyer annuel", "110 K€", INK),
                 ("% du CA", "31 %", INK),
                 ("EBE An 1", "+38 K€", GREEN),
                 ("Cash net An 1", "+32 K€", GREEN)],
                right_accent=GREEN)
    add_footer(s, "Acte 3 · Argent")

    # ---------- 09 Scénarios ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Scénarios An 1")
    add_pagination(s, 9, 13)
    add_title(s, CONTENT_TOP, [
        ("Trois trajectoires,", False, True),
        ("une seule vivable", True, False),
        (".", False, False),
    ], size=44)
    add_accent_line(s, Inches(3.0))
    y = Inches(3.3)
    y = add_step(s, y, 1, "Pessimiste",
                 "Bail non renégocié (180 K€) · CA stable à 350 K€",
                 amount="−8 K€", amount_color=RED)
    y = add_step(s, y, 2, "Cible",
                 "Bail à 110 K€ · CA stable à 350 K€ · effectif inchangé",
                 amount="+32 K€", amount_color=GREEN)
    y = add_step(s, y, 3, "Relance",
                 "Bail à 110 K€ · CA porté à 450 K€ · marketing local + offre",
                 amount="+89 K€", amount_color=GREEN)
    add_footer(s, "Acte 4 · Projection")

    # ---------- 10 Périmètre ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Périmètre")
    add_pagination(s, 10, 13)
    add_title(s, CONTENT_TOP, [
        ("Ce que le repreneur prend, ", False, True),
        ("ce qu'il laisse", True, False),
        (".", False, False),
    ], size=42)
    add_accent_line(s, Inches(3.0))
    add_two_col(s, Inches(3.3),
                "Repris",
                [("Fonds de commerce & clientèle", "✓", GREEN),
                 ("Bail commercial (renégocié)", "✓", GREEN),
                 ("Matériel & équipements", "inv. CP", INK),
                 ("Stocks (forfait)", "✓", GREEN),
                 ("Salariés repris (sur effectif)", "[N] / [M]", INK),
                 ("Contrats essentiels (assurance)", "L. 642-7", INK)],
                "Exclu",
                [("Créances clients antérieures", "✗", RED),
                 ("Dettes fournisseurs & fiscales", "✗", RED),
                 ("Congés payés antérieurs", "✗", RED),
                 ("Trésorerie & comptes bancaires", "✗", RED),
                 ("Contrats hors L. 642-7", "✗", RED),
                 ("Réserve de propriété & crédit-bail", "✗", RED)],
                right_accent=RED)
    add_footer(s, "Acte 4 · Projection")

    # ---------- 11 Ticket ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_eyebrow(s, "Ticket")
    add_pagination(s, 11, 13)
    add_title(s, CONTENT_TOP, [
        ("Cash à mobiliser :", False, True),
        ("100 à 125 K€", True, False),
        (".", False, False),
    ], size=44)
    add_accent_line(s, Inches(3.0))
    y = Inches(3.3)
    y = add_step(s, y, 1, "Prix d'offre au tribunal",
                 "8,6 % du CA de référence — conforme à la pratique distressed (≤ 12 %)",
                 amount="30 K€")
    y = add_step(s, y, 2, "CAPEX rafraîchissement",
                 "Sols, peintures, matériel à remettre à niveau — cible 6 premiers mois",
                 amount="35 K€")
    y = add_step(s, y, 3, "BFR de démarrage",
                 "3 mois de charges fixes — abonnements clients lissés",
                 amount="35 K€")
    y = add_step(s, y, 4, "Frais (acte, garantie, divers)",
                 "Honoraires rédacteur, mainlevée sûretés, garantie bancaire",
                 amount="15 K€")
    # total row
    line = s.shapes.add_connector(1, LEFT, y, SLIDE_W - RIGHT, y)
    line.line.color.rgb = INK
    line.line.width = Pt(1.5)
    tb, tf = add_textbox(s, LEFT, y + Inches(0.1), Inches(7), Inches(0.5))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Total cash mobilisé"
    style_run(run, font=SANS, size=18, bold=True, color=INK)
    tb, tf = add_textbox(s, SLIDE_W - RIGHT - Inches(3),
                         y + Inches(0.1), Inches(3), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "115 K€"
    style_run(run, font=MONO, size=22, bold=True, color=INK)
    add_footer(s, "Acte 4 · Projection")

    # ---------- 12 Conditions (DARK) ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, INK)
    add_eyebrow(s, "Conditions", dark=True)
    add_pagination(s, 12, 13, dark=True)
    # title in cream
    tb, tf = add_textbox(s, LEFT, CONTENT_TOP, SLIDE_W - LEFT - RIGHT, Inches(2))
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    r = p.add_run(); r.text = "Trois feux verts,"
    style_run(r, font=SANS, size=44, color=CREAM)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.05
    r = p2.add_run(); r.text = "avant signature"
    style_run(r, font=SERIF, size=44, italic=True, color=CREAM)
    r = p2.add_run(); r.text = "."
    style_run(r, font=SANS, size=44, color=CREAM)
    # accent line in cream
    line = s.shapes.add_connector(1, LEFT, Inches(3.0), LEFT + Inches(0.6), Inches(3.0))
    line.line.color.rgb = CREAM
    line.line.width = Pt(1.5)
    y = Inches(3.4)
    y = add_step(s, y, 1, "Bail renégocié ≤ 110 K€/an",
                 "Protocole tripartite repreneur / bailleresse / AJ — avant dépôt",
                 qualifier="Bloquant", qualifier_color=RGBColor(0xE5, 0x80, 0x78),
                 dark=True)
    y = add_step(s, y, 2, "Lettre d'intention bancaire 80 K€",
                 "Garantie de financement à joindre au pli — chèque de banque ou GAPD",
                 qualifier="Bloquant", qualifier_color=RGBColor(0xE5, 0x80, 0x78),
                 dark=True)
    y = add_step(s, y, 3, "Inventaire matériel contradictoire",
                 "Vérification physique avant offre — écart toléré ≤ 10 %",
                 qualifier="Confirmatif", dark=True)
    add_footer(s, "Acte 5 · Décision", dark=True)

    # ---------- 13 Recommandation (DARK) ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, INK)
    add_eyebrow(s, "Recommandation", dark=True)
    add_pagination(s, 13, 13, dark=True)
    # main title centered-ish
    tb, tf = add_textbox(s, LEFT, Inches(2.4), SLIDE_W - LEFT - RIGHT, Inches(2.5))
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    r = p.add_run(); r.text = "Déposer une offre à "
    style_run(r, font=SANS, size=54, color=CREAM)
    r = p.add_run(); r.text = "30 K€"
    style_run(r, font=SERIF, size=54, italic=True, color=CREAM)
    r = p.add_run(); r.text = ","
    style_run(r, font=SANS, size=54, color=CREAM)
    p2 = tf.add_paragraph(); p2.line_spacing = 1.05
    r = p2.add_run(); r.text = "conditionnelle à la renégociation du bail."
    style_run(r, font=SANS, size=54, color=CREAM)
    # subtitle
    tb, tf = add_textbox(s, LEFT, Inches(4.7), Inches(11), Inches(1.2))
    p = tf.paragraphs[0]; p.line_spacing = 1.5
    r = p.add_run()
    r.text = ("Ticket repreneur 115 K€ cash. Cash net An 1 attendu +32 K€ "
              "en scénario cible, +89 K€ si CA relancé. "
              "Trois conditions suspensives à valider avant dépôt.")
    style_run(r, font=SANS, size=15, color=RGBColor(0xC8, 0xC7, 0xC2))
    # CTA
    tb, tf = add_textbox(s, LEFT, Inches(5.85), SLIDE_W - LEFT - RIGHT, Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Décision interne avant le 14/05/2026."
    style_run(r, font=SANS, size=15, color=RGBColor(0x90, 0x90, 0x8E))
    tb, tf = add_textbox(s, LEFT, Inches(6.25), SLIDE_W - LEFT - RIGHT, Inches(0.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.2
    r = p.add_run(); r.text = "Dépôt greffe Nanterre — "
    style_run(r, font=SANS, size=22, color=CREAM)
    r = p.add_run(); r.text = "21/05/2026 · 12h"
    style_run(r, font=SERIF, size=22, italic=True, color=CREAM)
    r = p.add_run(); r.text = "."
    style_run(r, font=SANS, size=22, color=CREAM)
    add_footer(s, "brantham.fr", dark=True, url="brantham.fr")

    out = Path(__file__).parent / "pitch-repreneur.pptx"
    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
