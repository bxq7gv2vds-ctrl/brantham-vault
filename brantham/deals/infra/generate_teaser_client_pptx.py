from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("INFRA-teaser-client.pptx")

INK = RGBColor(16, 16, 15)
T2 = RGBColor(74, 74, 71)
T3 = RGBColor(125, 125, 120)
LINE = RGBColor(185, 184, 178)
OFF = RGBColor(246, 245, 242)
WHITE = RGBColor(255, 255, 255)

FONT = "Cambria"


def run_style(run, size=8, bold=False, italic=False, color=INK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def textbox(slide, x, y, w, h, text="", size=8, bold=False, italic=False, color=INK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 1.02
    r = p.add_run()
    r.text = text
    run_style(r, size=size, bold=bold, italic=italic, color=color)
    return shape


def multi_para(slide, x, y, w, h, parts, size=7.2):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, item in enumerate(parts):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        p.line_spacing = 1.02
        if isinstance(item, tuple):
            lead, body = item
            r = p.add_run()
            r.text = lead
            run_style(r, size=size, bold=True)
            r = p.add_run()
            r.text = " " + body
            run_style(r, size=size)
        else:
            r = p.add_run()
            r.text = item
            run_style(r, size=size)
    return shape


def hline(slide, x, y, w, color=LINE, width=0.7):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(width))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    shp.line.width = Pt(0)
    return shp


def vline(slide, x, y, h, color=LINE, width=0.7):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Pt(width), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    shp.line.width = Pt(0)
    return shp


def band(slide, x, y, w, h, fill=OFF):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    shp.line.width = Pt(0)
    return shp


def header(slide, kicker, title, subtitle, page):
    textbox(slide, 0.42, 0.22, 1.9, 0.18, kicker.upper(), size=5.8, color=T3)
    textbox(slide, 0.42, 0.43, 6.2, 0.36, title, size=15.5, bold=True)
    textbox(slide, 0.42, 0.82, 7.6, 0.20, subtitle, size=6.7, color=T2)
    textbox(slide, 8.15, 0.28, 1.35, 0.16, "Brantham Partners", size=6.8, bold=True, align=PP_ALIGN.RIGHT)
    textbox(slide, 8.15, 0.50, 1.35, 0.16, "pre-NDA confidentiel", size=5.7, color=T3, align=PP_ALIGN.RIGHT)
    hline(slide, 0.42, 1.08, 9.10, INK, 0.85)
    textbox(slide, 9.23, 7.12, 0.28, 0.12, str(page), size=5.6, color=T3, align=PP_ALIGN.RIGHT)
    textbox(slide, 0.42, 7.12, 4.2, 0.12, "Document confidentiel — diffusion interdite sans accord ecrit prealable", size=5.4, color=T3)


def mini_table(slide, x, y, col_widths, row_h, headers, rows, sizes=None):
    total_w = sum(col_widths)
    band(slide, x, y, total_w, row_h, WHITE)
    hline(slide, x, y, total_w, INK, 0.9)
    hline(slide, x, y + row_h, total_w, LINE, 0.6)
    cx = x
    for i, head in enumerate(headers):
        textbox(slide, cx + 0.03, y + 0.06, col_widths[i] - 0.06, row_h - 0.05, head, size=(sizes or {}).get("header", 5.8), bold=True)
        cx += col_widths[i]
    cy = y + row_h
    for r_idx, row in enumerate(rows):
        if r_idx % 2 == 1:
            band(slide, x, cy, total_w, row_h, OFF)
        cx = x
        for c_idx, val in enumerate(row):
            align = PP_ALIGN.RIGHT if c_idx > 0 and str(val).strip()[:1] in "~+-0123456789" else PP_ALIGN.LEFT
            textbox(slide, cx + 0.03, cy + 0.055, col_widths[c_idx] - 0.06, row_h - 0.05, str(val), size=(sizes or {}).get("body", 5.95), color=INK, align=align)
            cx += col_widths[c_idx]
        hline(slide, x, cy + row_h, total_w, LINE, 0.35)
        cy += row_h
    hline(slide, x, cy, total_w, INK, 0.8)
    return cy


def kpi_grid(slide, x, y, w, h, items):
    cols = len(items)
    cell_w = w / cols
    hline(slide, x, y, w, INK, 0.85)
    hline(slide, x, y + h, w, INK, 0.85)
    for i, (label, value, note) in enumerate(items):
        cx = x + i * cell_w
        if i:
            vline(slide, cx, y + 0.05, h - 0.10, LINE, 0.5)
        textbox(slide, cx + 0.06, y + 0.09, cell_w - 0.12, 0.13, label.upper(), size=4.8, color=T3, align=PP_ALIGN.CENTER)
        textbox(slide, cx + 0.06, y + 0.30, cell_w - 0.12, 0.28, value, size=11.5, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, cx + 0.06, y + 0.62, cell_w - 0.12, 0.20, note, size=5.0, color=T2, align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# Slide 1
s = prs.slides.add_slide(blank)
s.background.fill.solid()
s.background.fill.fore_color.rgb = WHITE
header(
    s,
    "Opportunite de reprise",
    "Agence de communication et production visuelle",
    "Marque historique du Grand Est — pres de 50 ans d'existence — procedure collective, statut de cession a confirmer",
    1,
)

kpi_grid(
    s,
    0.42,
    1.28,
    9.10,
    0.92,
    [
        ("Anciennete", "≈50 ans", "marque installee"),
        ("CA pic", "3,0 M€", "2021/22"),
        ("CA dernier plein", "2,3 M€", "31/03/2025"),
        ("CA recent", "1,8 M€", "annualise 9m"),
        ("Effectif", "18-19", "salaries actuels"),
        ("Point mort", "2,4-2,6 M€", "CA estime"),
    ],
)

textbox(s, 0.42, 2.42, 1.9, 0.16, "TRAJECTOIRE FINANCIERE", size=5.7, color=T3, bold=True)
mini_table(
    s,
    0.42,
    2.65,
    [1.15, 0.76, 0.76, 0.76, 0.76, 0.90],
    0.255,
    ["K€", "03/22", "03/23", "03/24", "03/25", "12/25 9m"],
    [
        ["CA", "2 993", "2 427", "2 216", "2 312", "1 357"],
        ["Marge brute", "2 208", "1 854", "1 738", "1 718", "1 043"],
        ["Marge brute %", "73,8%", "76,4%", "78,4%", "74,3%", "76,9%"],
        ["Salaires + charges", "1 137", "1 064", "961", "972", "820"],
        ["EBE", "+227", "+36", "+22", "+79", "-284"],
        ["Resultat net", "+97", "-77", "-73", "-1", "-345"],
    ],
    {"header": 5.2, "body": 5.25},
)
multi_para(
    s,
    5.75,
    2.60,
    3.77,
    1.70,
    [
        ("Lecture.", "Le CA recent annualise ressort autour de 1,8 M€, soit environ -22% vs dernier exercice plein et environ -40% vs le pic recent."),
        ("Marge preservee.", "La marge brute reste proche de 75-77% : le probleme n'est pas le mix vendu, mais le volume face aux couts fixes."),
        ("Point mort.", "Le seuil EBE est estime a 2,4-2,6 M€ de CA : l'agence est 600-800 K€ sous son niveau d'equilibre."),
    ],
    size=6.25,
)

textbox(s, 0.42, 4.72, 1.75, 0.16, "CAUSES DCP", size=5.7, color=T3, bold=True)
mini_table(
    s,
    0.42,
    4.95,
    [1.45, 2.45, 2.35],
    0.255,
    ["Cause", "Fait constate", "Impact economique"],
    [
        ["Dette LBO INTRA", "Rachat 2016 via holding INTRA", "Remontees de cash / stress financier"],
        ["Conjoncture", "Budgets communication coupes en priorite", "CA en contraction, decisions clients allongees"],
        ["Choc IA", "Clients persuades de produire eux-memes", "Pression sur prestations simples"],
        ["Absence commerciale", "Dirigeante commerciale absente > 1 an", "Pipeline affaibli, traction commerciale perdue"],
        ["Couts fixes", "Masse salariale inerte face au CA", "EBE negatif malgre marge brute correcte"],
        ["Tresorerie", "VMP consommees, cash proche zero", "Sauvegarde avant cessation complete"],
    ],
    {"header": 5.05, "body": 5.05},
)
multi_para(
    s,
    7.05,
    4.94,
    2.47,
    1.60,
    [
        ("Synthese.", "Le retournement resulte d'un effet ciseaux : baisse rapide du chiffre d'affaires, structure de couts encore dimensionnee pour un volume superieur, dette de LBO portee par la holding et choc IA sur les prestations commoditisees."),
    ],
    size=6.15,
)

# Slide 2
s = prs.slides.add_slide(blank)
s.background.fill.solid()
s.background.fill.fore_color.rgb = WHITE
header(
    s,
    "These repreneur",
    "Actifs, leviers de reprise et points a securiser",
    "Base exploitable, mais offre a conditionner a la transferabilite des contrats, au bail actuel et au statut procedural",
    2,
)

textbox(s, 0.42, 1.27, 1.7, 0.16, "CE QUI A DE LA VALEUR", size=5.7, color=T3, bold=True)
mini_table(
    s,
    0.42,
    1.50,
    [1.70, 3.15],
    0.255,
    ["Actif", "Lecture repreneur"],
    [
        ["Marque historique", "Pres de 50 ans d'existence, credibilite locale difficile a reconstruire"],
        ["Ancrage Grand Est", "Proximite PME, ETI industrielles, institutions et comptes regionaux"],
        ["Production integree", "Photo, video, redaction, print, digital ; moins de dependance a la sous-traitance"],
        ["Positionnement responsable", "Print eco-responsable, digital frugal, Green UX ; angle RSE exploitable"],
        ["Portefeuille B2B", "Cas clients recrutement, produit, salon, institutionnel ; valeur a confirmer par contrats"],
        ["Equipe resserree", "18-19 salaries ; premiere reduction deja engagee vs 21 en 2023 / 25 plaquette"],
        ["Incorporels", "Marque, methodologie, phototheque/videotheque, noms de domaine, fichier clients"],
    ],
    {"header": 5.15, "body": 5.05},
)

textbox(s, 5.55, 1.27, 1.7, 0.16, "THESE CHIFFREE", size=5.7, color=T3, bold=True)
mini_table(
    s,
    5.55,
    1.50,
    [1.45, 1.35, 1.17],
    0.255,
    ["Levier", "Base", "Effet"],
    [
        ["Retour point mort", "2,4-2,6 M€ CA", "EBE positif"],
        ["Relance commerciale", "CA historique 3,0 M€", "Volume recuperable"],
        ["Recalibrage social", "18-19 salaries", "Cout fixe ajuste"],
        ["Dette historique", "PGE ~850 K€", "Non repris"],
        ["Charge L.642-12", "~9,8 K€", "Impact limite"],
        ["Materiel revendique", "10-15 K€", "Negociable"],
        ["BFR redemarrage", "100-200 K€", "A financer"],
    ],
    {"header": 5.15, "body": 5.05},
)

textbox(s, 0.42, 4.06, 1.85, 0.16, "ACTIONS DEJA ENGAGEES", size=5.7, color=T3, bold=True)
mini_table(
    s,
    0.42,
    4.29,
    [1.95, 2.95],
    0.255,
    ["Action", "Lecture"],
    [
        ["Plan d'economies 2023-25", "Effort de restructuration deja amorce"],
        ["Effectif 21 -> 18/19", "Base sociale deja resserree par non-remplacements"],
        ["Sous-traitance variabilisee", "Modele potentiellement plus flexible"],
        ["Primes direction supprimees", "Charges non operationnelles reduites"],
        ["Renforcement commercial", "Tentative de relance a reprendre / challenger"],
    ],
    {"header": 5.1, "body": 5.0},
)

textbox(s, 5.55, 4.06, 1.85, 0.16, "POINTS A CONFIRMER", size=5.7, color=T3, bold=True)
mini_table(
    s,
    5.55,
    4.29,
    [1.70, 2.27],
    0.255,
    ["Sujet", "Pourquoi c'est critique"],
    [
        ["Statut procedural", "Sauvegarde : cession effective a clarifier"],
        ["Contrats clients", "Transferabilite / concentration / CA recurrent"],
        ["Bail Strasbourg", "Condition d'exploitation post-reprise"],
        ["Etat du passif", "Privileges et charges residuelles"],
        ["Effectif actualise", "Volet social et conges payes repris"],
        ["PI / domaines / fichier", "Valeur incorporelle a securiser"],
    ],
    {"header": 5.1, "body": 5.0},
)

multi_para(
    s,
    0.42,
    6.18,
    6.15,
    0.55,
    [
        ("Conclusion.", "Le dossier n'est pas a vendre comme une agence saine, mais comme une plateforme regionale a relancer : marque ancienne, savoir-faire integre, marge brute preservee, dette historique non reprise et point mort clairement identifiable."),
    ],
    size=6.25,
)
textbox(
    s,
    6.90,
    6.17,
    2.62,
    0.64,
    "Disclaimer. Presentation preliminaire et confidentielle ; ne constitue ni une offre ni un memorandum d'information. Informations d'identification volontairement omises a ce stade. Contact : Paul Roulleau — Brantham Partners.",
    size=5.6,
    color=T2,
)

prs.save(OUT)
print(OUT)
